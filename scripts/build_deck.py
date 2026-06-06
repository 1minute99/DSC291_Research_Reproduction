"""Declarative builder for the MILAN reproduction deck.

Builds the full deck from the template's slide master/layouts (so the footer,
slide numbers and theme are preserved), following the project template's
recommended page counts and the example deck's multi-slide-with-figures format:

  Related works ~2 · Key ideas ~3 · Key results ~2 · Reproduction ~3 ·
  Additional experiments 3 (one per extension) · + intro/discussion/refs.

All numbers are the 10% overlay run. Figures are reused from results/figs/;
comparison tables are native PowerPoint tables. Output -> the MILAN pptx.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "SP 26 DSC 291 Project Template - Research Reproduction.pptx"
OUT = ROOT / "SP 26 DSC 291 MILAN - Research Reproduction .pptx"

FIG = ROOT / "results" / "figs"
SLIDE = FIG / "slides"
PAPER = FIG / "paper"

# figure shortcuts
PIPELINE = SLIDE / "milan_pipeline_diagram.png"
GRID = SLIDE / "spurious_dataset_grid.png"
ARCHBAR = SLIDE / "arch_text_neuron_bar.png"
SUMMARY = SLIDE / "summary_metrics_table.png"
FIG7 = FIG / "fig7.png"
FIG8 = FIG / "fig8.png"
ARCH = FIG / "fig_arch_comparison.png"
LAYER = FIG / "fig_layer_analysis.png"
CLIP = FIG / "fig_clip_analysis.png"
PFIG7 = PAPER / "paper_fig7.png"
PFIG8 = PAPER / "paper_fig8.png"

BULLET = "●"
SUBBULLET = "◦"

# slide geometry (EMU). Title box is tall enough for a 2-line title at 24pt;
# the body starts below it so wrapped titles never overlap the content.
TITLE_BOX = (311700, 360000, 8520600, 840000)
TITLE_PT = 24
BODY_LEFT, BODY_TOP, BODY_WIDTH = 311700, 1_300_000, 8520600
BODY_BOTTOM = 4858000  # leave room above the footer

DARK = RGBColor(0x2E, 0x3B, 0x4E)
LIGHTROW = RGBColor(0xEE, 0xF2, 0xF7)
WHITEROW = RGBColor(0xFB, 0xFC, 0xFE)


# ---------------------------------------------------------------- paragraph fmt
def _para_fmt(paragraph, kind):
    p = paragraph._p
    old = p.find(qn("a:pPr"))
    if old is not None:
        p.remove(old)
    if kind == "bullet":
        xml = (f'<a:pPr {nsdecls("a")} marL="285750" indent="-285750" lvl="0" algn="l">'
               '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="600"/></a:spcAft>'
               f'<a:buSzPts val="1500"/><a:buChar char="{BULLET}"/></a:pPr>')
    elif kind == "item":
        xml = (f'<a:pPr {nsdecls("a")} marL="742950" indent="-285750" lvl="1" algn="l">'
               '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="300"/></a:spcAft>'
               f'<a:buSzPts val="1300"/><a:buChar char="{SUBBULLET}"/></a:pPr>')
    elif kind == "header":
        xml = (f'<a:pPr {nsdecls("a")} marL="0" indent="0" lvl="0" algn="l">'
               '<a:spcBef><a:spcPts val="600"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="300"/></a:spcAft>'
               '<a:buNone/></a:pPr>')
    else:
        return
    p.insert(0, parse_xml(xml))


def _fill_body(ph, lines, font_pt):
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, entry in enumerate(lines):
        text, kind = entry if isinstance(entry, tuple) else (entry, "bullet")
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        _para_fmt(para, kind)
        for run in para.runs:
            run.font.size = Pt(font_pt)
            if kind == "header":
                run.font.bold = True


def _set_title(slide, text, font_pt=None):
    ph = slide.placeholders[0]
    ph.left, ph.top, ph.width, ph.height = (Emu(v) for v in TITLE_BOX)
    tf = ph.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    tf.paragraphs[0].text = text
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(font_pt or TITLE_PT)


def place_fig(slide, path, region, valign="center"):
    """Fit an image (preserving aspect) inside region=(left,top,w,h)."""
    l, t, w, h = region
    im = Image.open(path)
    iw, ih = im.size
    im.close()
    tw = w
    th = int(tw * ih / iw)
    if th > h:
        th = h
        tw = int(th * iw / ih)
    x = l + (w - tw) // 2
    y = t if valign == "top" else (t + (h - th) // 2 if valign == "center" else t + (h - th))
    slide.shapes.add_picture(str(path), x, y, width=tw, height=th)


def place_row(slide, paths, region, gap=120000, valign="center"):
    l, t, w, h = region
    aspects = []
    for p in paths:
        im = Image.open(p)
        aspects.append(im.width / im.height)
        im.close()
    n = len(paths)
    avail = w - gap * (n - 1)
    H = min(h, int(avail / sum(aspects)))
    widths = [int(H * a) for a in aspects]
    roww = sum(widths) + gap * (n - 1)
    x = l + (w - roww) // 2
    y = t if valign == "top" else t + (h - H) // 2
    for p, ww in zip(paths, widths):
        slide.shapes.add_picture(str(p), x, y, width=ww, height=H)
        x += ww + gap


def add_table(slide, rows, region, col_fracs=None, font_pt=11, header=True):
    l, t, w, h = region
    nrows, ncols = len(rows), len(rows[0])
    gframe = slide.shapes.add_table(nrows, ncols, Emu(l), Emu(t), Emu(w), Emu(h))
    tbl = gframe.table
    if col_fracs:
        for j, fr in enumerate(col_fracs):
            tbl.columns[j].width = Emu(int(w * fr))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Emu(18000)
            cell.margin_bottom = Emu(18000)
            cell.margin_left = Emu(55000)
            cell.margin_right = Emu(55000)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(font_pt)
                if i == 0 and header:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                elif j == 0:
                    run.font.bold = True
            if i == 0 and header:
                cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHTROW if i % 2 else WHITEROW
    return gframe


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- build helpers
def build():
    prs = Presentation(str(TEMPLATE))
    LAY_TITLE = prs.slide_layouts[0]
    LAY_SECTION = prs.slide_layouts[1]
    LAY_BODY = prs.slide_layouts[2]

    # wipe the template's slides; keep master/layouts. Drop the presentation->
    # slide relationship too (not just the sldId), so the old parts are pruned
    # on save instead of lingering and colliding with new slide part names.
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        prs.part.drop_rel(sld.get(qn("r:id")))
        xml_slides.remove(sld)

    def section(title):
        s = prs.slides.add_slide(LAY_SECTION)
        s.placeholders[0].text = title
        return s

    def title_slide(title, subtitle_lines, title_pt=36, sub_pt=15):
        s = prs.slides.add_slide(LAY_TITLE)
        t0 = s.placeholders[0].text_frame
        t0.word_wrap = True
        t0.text = title
        for run in t0.paragraphs[0].runs:
            run.font.size = Pt(title_pt)
        tf = s.placeholders[1].text_frame
        tf.word_wrap = True
        tf.clear()
        for i, ln in enumerate(subtitle_lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = ln
            for run in para.runs:
                run.font.size = Pt(sub_pt)
        return s

    def content(title, lines=None, body_h=None, font_pt=15, title_pt=None,
                figs=None, fig_region=None, fig_row=None, table=None,
                notes=None):
        """figs: list of (path, region). fig_region/fig_row: convenience for a
        single fig / a row placed in the area below the body."""
        s = prs.slides.add_slide(LAY_BODY)
        _set_title(s, title, title_pt)
        body = s.placeholders[1]
        if lines:
            bh = body_h if body_h else (BODY_BOTTOM - BODY_TOP)
            body.left, body.top = Emu(BODY_LEFT), Emu(BODY_TOP)
            body.width, body.height = Emu(BODY_WIDTH), Emu(bh)
            _fill_body(body, lines, font_pt)
            area_top = BODY_TOP + bh + 80000
        else:
            # remove the empty body placeholder so it doesn't show a prompt
            body._element.getparent().remove(body._element)
            area_top = BODY_TOP
        area = (BODY_LEFT, area_top, BODY_WIDTH, BODY_BOTTOM - area_top)
        if fig_region is not None:
            auto = fig_region[1] is None
            reg = area if auto else fig_region[1]
            # auto-placed figures sit at the bottom so a wide/short figure can't
            # float up into the bullets above it
            va = "bottom" if auto else (fig_region[2] if len(fig_region) > 2 else "center")
            place_fig(s, fig_region[0], reg, valign=va)
        if fig_row is not None:
            place_row(s, fig_row, area)
        if figs:
            for path, region in figs:
                place_fig(s, path, region)
        if table is not None:
            add_table(s, table["rows"], table.get("region", area),
                      col_fracs=table.get("col_fracs"), font_pt=table.get("font", 11))
        if notes:
            set_notes(s, notes)
        return s

    half = BODY_WIDTH // 2 - 80000

    # =========================================================== TITLE / CITATION
    title_slide("MILAN — Editing Spurious Features",
                ["Reproducing Hernandez et al. (ICLR 2022) — Natural Language "
                 "Descriptions of Deep Visual Features",
                 "+ extensions: VGG16 · ResNet18 layer-depth · CLIP ViT-B/32",
                 "Group [#] · Wonmin Kim · Seongho Kim · Ming-Yang Wu · Steven Tsai · June 9, 2026"])

    title_slide("Paper [A]: Natural Language Descriptions of Deep Visual Features",
                ["[A] Hernandez, Schwettmann, Bau, Bagashvili, Oliva, Andreas. ICLR 2022.",
                 "arXiv:2201.11114 · project: milan.csail.mit.edu",
                 "Focus: Section 7 — “Editing Spurious Features”"])

    # =========================================================== SEC 1
    section("Sec 1: Introduction & Motivation")

    content("Overview: interpretability for trustworthy ML",
            ["Deep nets are accurate but opaque — we can't easily tell what each neuron encodes.",
             "Prior interpretability (Network Dissection, TCAV, CLIP-Dissect) is tied to a fixed concept vocabulary.",
             "MILAN goes open-vocabulary: a learned image-conditioned language model writes the description.",
             "Section 7 payoff: use those descriptions to find and remove neurons that latched onto a spurious shortcut."],
            body_h=1_950_000, font_pt=15,
            fig_region=(PIPELINE, (BODY_LEFT, 3_520_000, BODY_WIDTH, 1_260_000), "top"),
            notes="45s. Earlier methods need a concept list up front; MILAN learns captions from 60k human "
                  "annotations and produces free-form descriptions, which become actionable handles for editing.")

    content("The problem: spurious shortcuts inside neurons",
            ["A model trained on data with a spurious correlation latches onto the shortcut — here, the class "
             "name painted in the image corner.",
             "Result: high clean accuracy, but adversarial accuracy collapses when the corner text is wrong.",
             "Open problem: identify which neurons encode the shortcut, at scale, without retraining."],
            body_h=1_500_000, font_pt=15,
            fig_region=(GRID, None, "center"),
            notes="Top row clean, middle row training set with class-name overlay, bottom row adversarial "
                  "(wrong-class) overlay. Reading the text instead of the object is the failure mode.")

    content("Why this matters for trustworthy ML",
            ["Robustness: text overlays, watermarks and background bias silently degrade real-world deployment.",
             "Fairness & auditing: knowing what a model relies on is a prerequisite for accountability.",
             "Editing without retraining: full retraining is expensive and brittle — neuron-level edits are surgical and reversible.",
             "Three course themes — spurious correlations, adversarial robustness, model editing — meet in one self-contained loop."],
            font_pt=16,
            notes="Tie back to the course: this single experiment exercises spurious correlations, adversarial "
                  "robustness, and interpretability-driven editing.")

    # =========================================================== SEC 2
    section("Sec 2: Related works")

    content("Related work: fixed lists → open-vocabulary captions",
            ["Every prior method requires you to enumerate concepts up front — so anything off-vocabulary "
             "(like “white text in the corner”) is missed unless anticipated.",
             "MILAN removes that bottleneck by learning to caption neurons in free-form English."],
            body_h=1_150_000, font_pt=14,
            table={"rows": [
                ["Method", "Concept source", "Output", "Limitation"],
                ["Network Dissection (2017)", "~1,200 fixed labels", "best-matching concept", "closed vocabulary"],
                ["TCAV (2018)", "user concept examples", "directional sensitivity", "needs probe sets per concept"],
                ["CLIP-Dissect (2023)", "CLIP text encoder", "nearest CLIP concept", "still vocab-bound"],
                ["MILAN (this paper)", "60k human captions", "free-form description", "harder to evaluate"],
            ], "col_fracs": [0.26, 0.24, 0.24, 0.26], "font": 11.5},
            notes="60s. Walk the family tree. Common thread: prior work needs a concept list; MILAN's vocabulary "
                  "is whatever appears in 60k human neuron captions.")

    content("Where MILAN fits — and the trade-off",
            ["Open-vocabulary captions are far more expressive, but harder to score objectively.",
             "The paper's fix is a PMI decoding objective (next section) that rewards specificity and penalises generic captions.",
             "For editing, the key advantage is that descriptions are directly actionable: regex-flag captions mentioning text, then ablate.",
             "Downstream filter is the weak link — a string match for 'text'/'word'/'letter' (we return to this in Limitations)."],
            font_pt=16,
            notes="Position MILAN as the open-vocab successor; foreshadow the PMI objective and the regex-filter limitation.")

    # =========================================================== SEC 3
    section("Sec 3: Main methods")

    content("Key idea 1 — neurons as top-activating exemplars",
            ["Each conv channel has a set of image patches that maximally activate it — its exemplars.",
             "MILAN feeds those exemplars into a small image-conditioned encoder-decoder language model.",
             "The decoder is trained on MILANNOTATIONS: 60k human descriptions of such exemplar sets.",
             "Output: a free-form English caption per neuron, e.g. “words and letters” or “dog faces”."],
            body_h=2_050_000, font_pt=15,
            fig_region=(PIPELINE, (BODY_LEFT, 3_520_000, BODY_WIDTH, 1_260_000), "top"),
            notes="The pipeline: exemplars -> image-conditioned LM -> caption. Trained on 60k human exemplar captions.")

    content("Key idea 2 — the PMI decoding objective",
            [("What it solves", "header"),
             ("Plain likelihood decoding produces bland captions ('an image', 'a photo').", "item"),
             ("MILAN instead maximises pointwise mutual information:", "item"),
             ("    argmax_d  log p(d | exemplars)  −  λ·log p(d)", "item"),
             ("The −log p(d) term penalises descriptions that are likely regardless of the neuron, pushing toward specificity.", "item"),
             ("Why it matters for us", "header"),
             ("Specific captions ('words', 'letters') are exactly what makes text-selective neurons findable by a keyword filter.", "item")],
            font_pt=15,
            notes="PMI = log p(d|exemplars) - log p(d). The prior term kills generic captions; specificity is what "
                  "makes the text neurons identifiable downstream.")

    content("Key idea 3 — editing spurious features",
            ["1. Train a classifier on spurious-text data (class name painted in a corner).",
             "2. Extract top-k exemplars for every neuron; caption each with MILAN.",
             "3. Flag text-selective neurons by regex ('text' / 'word' / 'letter').",
             "4. Ablate (zero) the flagged neurons and re-measure adversarial accuracy.",
             "No retraining, no gradient surgery — the captions alone tell you what to remove."],
            body_h=2_100_000, font_pt=15,
            fig_region=(GRID, None, "center"),
            notes="The self-contained editing loop. Steps 1-5 form the experiment; we extend step 1 to VGG16/CLIP "
                  "and add a per-layer analysis at step 4.")

    # =========================================================== SEC 4
    section("Sec 4: Experiment")

    content("Experimental setup — dataset",
            ["Spurious-text Imagenette: 10 ImageNet classes, class name painted in the top-left corner.",
             "Training set: 10% of images carry their (correct) class-name overlay — a predictive shortcut.",
             "Adversarial test set: every image carries a wrong-class overlay (the robustness probe).",
             "We use Imagenette (a public 10-class ImageNet subset) in place of the paper's full-ImageNet generator (compute)."],
            body_h=1_700_000, font_pt=15,
            fig_region=(GRID, None, "center"),
            notes="Dataset substitution noted. 10% overlay is the canonical setting (it gives the cleanest editing "
                  "separation; see speaker notes on the reproduction slides).")

    content("Experimental setup — model, MILAN & metrics",
            [("Model", "header"),
             ("ResNet18 trained from scratch (Appendix-E hyperparameters); also VGG16 and zero-shot CLIP ViT-B/32 for the extensions.", "item"),
             ("MILAN", "header"),
             ("Pretrained 'base' decoder from milan.csail.mit.edu — the authors' exact checkpoint.", "item"),
             ("Metrics", "header"),
             ("Clean val accuracy; adversarial test accuracy; ablation curve (adv-acc vs # neurons zeroed); text-neuron count.", "item"),
             ("Compute", "header"),
             ("Originally DSMLP GTX 1080 Ti (11 GB) — several OOM patches; reruns on a local RTX 5090.", "item")],
            font_pt=14,
            notes="Baselines for the ablation curve: sort-all (importance-ordered) and random (5 trials). "
                  "Everything except the dataset base follows the paper.")

    content("Paper's key result 1 — MILAN finds the text neuron",
            ["The paper trains a 10-class ImageNet classifier with class-name text stamped in the corner.",
             "(a) training examples · (b) adversarial test set with wrong-class text.",
             "(c) the top exemplars of neuron layer3-134, which MILAN captions “words and letters” — exactly the spurious feature."],
            body_h=1_500_000, font_pt=15,
            fig_region=(PFIG7, None, "center"),
            notes="Fig 7 from the paper. MILAN's caption pinpoints the shortcut-bearing neuron in plain English.")

    content("Paper's key result 2 — ablation recovers robustness",
            ["Y-axis: adversarial accuracy; X-axis: number of neurons zeroed.",
             "Blue 'sort text' (ablate MILAN text neurons) climbs from ~58% to ~63%, beating the orange 'sort all' importance baseline.",
             "Dashed line: the no-text-distractor ceiling. Editing closes much of the gap without retraining."],
            body_h=1_500_000, font_pt=15,
            fig_region=(PFIG8, None, "center"),
            notes="Fig 8 from the paper — the headline editing result we set out to reproduce.")

    # ----- Code reproduction (3 slides) -----
    content("Code reproduction — setup & headline numbers",
            ["Same pipeline on spurious-text Imagenette (10% overlay), the authors' MILAN 'base' decoder, and the paper's ablation protocol.",
             "We reproduce the qualitative story; absolute numbers differ (Imagenette vs full ImageNet, 10 vs 1000 classes)."],
            body_h=1_150_000, font_pt=14,
            table={"rows": [
                ["Metric", "Paper (trend)", "Ours (10% overlay)"],
                ["Clean accuracy", "high", "75.4%"],
                ["Adversarial accuracy (no ablation)", "~58%", "51.5%"],
                ["MILAN text neurons", "a minority", "216 / 1024 (21.1%)"],
                ["Effect of ablating text neurons", "recovers adv-acc, beats baselines", "51.5% → ~55%, beats random/importance"],
            ], "col_fracs": [0.40, 0.30, 0.30], "font": 12},
            notes="Headline comparison. Trends match; magnitudes differ because of the dataset substitution.")

    content("Code reproduction — text neurons (our Fig 7)",
            ["Eight MILAN-flagged text neurons from our ResNet18, with the captions as row labels.",
             "The painted class words visible in the exemplar thumbnails ('pump', 'truck', 'springer'…) are exactly what these neurons fire on — MILAN is finding the right thing."],
            body_h=1_250_000, font_pt=14,
            fig_region=(FIG7, None, "center"),
            notes="Our reproduction of Fig 7. Captions like 'words'/'letters'; the overlaid words are what the neurons respond to.")

    content("Code reproduction — ablation curve (our Fig 8)",
            ["From the 51.5% no-ablation baseline, zeroing MILAN's text neurons (blue) recovers adversarial accuracy to ~55% and holds.",
             "Ablating the same count by importance (orange) or at random (gray) instead degrades to ~46% / ~37%.",
             "So caption-guided editing clearly beats both baselines — the paper's Section-7 claim reproduces."],
            body_h=1_500_000, font_pt=15,
            fig_region=(FIG8, None, "center"),
            notes="Our Fig 8 at 10% overlay. Text-sorted is the only curve that stays above the no-ablation line.")

    # ----- Additional experiments (3 slides) -----
    content("Additional experiment 1 — VGG16 generalization",
            ["VGG16 trained on the same 10% spurious set; MILAN flags 357/1472 (24.3%) text neurons.",
             "VGG16 is more shortcut-reliant than ResNet18 (adv baseline 19.8% vs 51.5%).",
             "Yet text-sorted ablation still recovers it strongly: 19.8% → 32.7% (+12.8pp), clearly above random — the editing generalizes across CNNs with no code changes."],
            body_h=1_700_000, font_pt=14,
            figs=[(ARCH, (BODY_LEFT, 2_950_000, half, BODY_BOTTOM - 2_950_000)),
                  (ARCHBAR, (BODY_LEFT + half + 160000, 2_950_000, half, BODY_BOTTOM - 2_950_000))],
            notes="Left: per-architecture ablation curves (solid=text-sorted, dashed=sort-all). Right: text-neuron "
                  "fraction across architectures. VGG16 relies on the shortcut more, but the edit still works.")

    content("Additional experiment 2 — where the shortcut lives",
            ["Per-layer MILAN text-neuron fraction in ResNet18: 9 → 27 → 37 → 77 → 66 (conv1 → layer4).",
             "It jumps 3.0× from conv1 (14%) to layer1 (42%) — the shortcut appears right after the first residual block — then thins out with depth.",
             "Description diversity falls monotonically with depth: deeper neurons are more redundant / object-like."],
            body_h=1_700_000, font_pt=14,
            fig_region=(LAYER, (BODY_LEFT, 2_950_000, BODY_WIDTH, BODY_BOTTOM - 2_950_000), "center"),
            notes="The shortcut crystallizes very early (layer1). Two panels: text-neuron fraction (left) and "
                  "caption diversity (right) by depth.")

    content("Additional experiment 3 — is CLIP more robust?",
            ["Run MILAN on CLIP's vision transformer blocks (no task training on the spurious set).",
             "Max text-neuron fraction across blocks is 7.9% (61/768) — about 2.7× lower than ResNet18 (21.1%).",
             "Top words are 'objects', 'colored', 'white' — not 'words'/'letters'. CLIP's vision-language pretraining is structurally more robust to text overlays."],
            body_h=1_700_000, font_pt=14,
            fig_region=(CLIP, (BODY_LEFT, 2_950_000, BODY_WIDTH, BODY_BOTTOM - 2_950_000), "center"),
            notes="CLIP zero-shot: far fewer text neurons, and they describe objects not text. Quantitative evidence "
                  "that vision-language pretraining resists the shortcut.")

    # =========================================================== SEC 5
    section("Sec 5: Concluding remarks")

    content("Conclusions & take-away",
            ["MILAN reproduces on Imagenette (10% overlay): 216/1024 text neurons; ablation recovers adv-acc 51.5% → ~55%, beating random/importance.",
             "Editing generalizes across CNN architectures (ResNet18 → VGG16) with no code changes.",
             "Spurious shortcuts crystallize right after the first residual block, then thin out with depth.",
             "CLIP's image-language pretraining is ~2.7× sparser in text-selective neurons than a task-trained CNN.",
             "Open-vocabulary, model-agnostic neuron descriptions are a practical handle for editing model behaviour."],
            body_h=2_350_000, font_pt=14,
            fig_region=(SUMMARY, (BODY_LEFT, 3_600_000, BODY_WIDTH, BODY_BOTTOM - 3_600_000), "center"),
            notes="Headline takeaways at 10% overlay, with the summary table of all four experiments.")

    # =========================================================== SEC 6
    section("Sec 6: Discussion")

    content("What we liked / didn't like",
            [("Liked", "header"),
             ("Open-vocabulary descriptions — no concept set to engineer.", "item"),
             ("Model-agnostic — same decoder on ResNet, VGG, CLIP.", "item"),
             ("Descriptions are actionable — they enable editing, not just analysis.", "item"),
             ("Didn't like", "header"),
             ("Text-neuron identification is a brittle regex ('caption', 'sign', 'character' slip through).", "item"),
             ("Decoder is English-only, bounded by the 60k-caption corpus.", "item"),
             ("No uncertainty on captions — confidently wrong looks like confidently right.", "item")],
            font_pt=14,
            notes="Honest take. Biggest like: escapes the concept-vocabulary bottleneck. Biggest dislike: the regex filter.")

    content("Strengths & weaknesses",
            [("Strengths", "header"),
             ("PMI objective is a clean, principled fix for generic captions.", "item"),
             ("Three case studies in the paper (BigGAN, ImageNet classifiers, adversarial editing) — real utility, not toy demos.", "item"),
             ("Pretrained decoder is plug-and-play across architectures.", "item"),
             ("Weaknesses", "header"),
             ("Decoder quality is bounded by the size/diversity of the 60k-caption corpus.", "item"),
             ("Transformer units need custom adapters — we wrote clip_glue.py to reshape ViT tokens to (B,C,H,W).", "item"),
             ("No reported test for caption consistency across seeds.", "item")],
            font_pt=14,
            notes="Structural strength: PMI. Structural weakness: inherits the caption corpus's blind spots.")

    content("Limitations of this paper",
            ["Keyword-based text-neuron identification is brittle — 'sign', 'character', 'label' all describe text but miss the regex.",
             "Ablation is destructive: a neuron is zeroed in full, with no graded / fine-grained edit.",
             "Only top-k exemplars reach the decoder — rare but discriminative activations may be invisible.",
             "Evaluation is ground-truth-free: no neuron-level gold labels, so caption correctness is largely qualitative.",
             "Our own finding: the editing benefit depends on shortcut sparsity — at 50% overlay the shortcut pervades >half the net and caption-guided ablation no longer separates from random."],
            font_pt=14,
            notes="The last bullet is our empirical contribution: editing needs an identifiable (sparse) shortcut. "
                  "Addressed-in-literature: embedding-based concept detection is an active area.")

    content("Potential future work",
            ["Embedding-based text-neuron detection: replace the regex with sentence-encoder similarity to catch synonyms/paraphrases.",
             "Graded edits: activation steering or low-rank shrinkage instead of zero-ablation — preserve capacity while removing the shortcut.",
             "Beyond top-k: condition the decoder on a fuller activation distribution to surface rare features.",
             "MILAN-for-transformers: extend to ViTs and LLM hidden states with shared adapters (start from our clip_glue.py).",
             "Harder spurious benchmarks: Waterbirds, CelebA, NICO — multi-axis spurious correlations beyond text overlay."],
            font_pt=14,
            notes="Embedding filter is the lowest-hanging fruit; activation-steering connects to recent mech-interp work.")

    content("References",
            ["[1] Hernandez, Schwettmann, Bau, Bagashvili, Oliva, Andreas. Natural Language Descriptions of Deep Visual Features. ICLR 2022. arXiv:2201.11114.",
             "[2] Bau, Zhou, Khosla, Oliva, Torralba. Network Dissection. CVPR 2017.",
             "[3] Kim et al. TCAV: Interpretability Beyond Feature Attribution. ICML 2018.",
             "[4] Oikarinen & Weng. CLIP-Dissect. ICLR 2023.",
             "[5] Radford et al. Learning Transferable Visual Models from Natural Language Supervision (CLIP). ICML 2021.",
             "[6] Howard. Imagenette — github.com/fastai/imagenette.",
             "[7] Our repo: github.com/1minute99/DSC291_Research_Reproduction."],
            font_pt=12)

    prs.save(str(OUT))
    print(f"saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
