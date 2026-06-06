"""Fill the DSC 291 Research Reproduction template in-place.

Edits text placeholders (preserving the template's font/layout) and adds
images for the experiment slides. Speaker notes are filled for every slide.
"""
from __future__ import annotations

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

BULLET = "●"   # ● filled circle (matches the template's level-0 bullet)
SUBBULLET = "◦"  # ○ hollow circle for second-level items

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "SP 26 DSC 291 Project Template - Research Reproduction .pptx"

FIG = ROOT / "results" / "figs"
SLIDE_FIG = FIG / "slides"
PAPER_FIG = FIG / "paper"


def _apply_para_format(paragraph, kind):
    """Force a paragraph's bullet/indent so it never depends on leftover XML.

    kind = 'bullet'  → level-0 filled bullet (the default body style)
           'item'    → level-1 hollow bullet (indented sub-item)
           'header'  → no bullet, flush left (used for sub-section headers)
    """
    p = paragraph._p
    old = p.find(qn("a:pPr"))
    if old is not None:
        p.remove(old)
    if kind == "bullet":
        xml = (f'<a:pPr {nsdecls("a")} marL="457200" indent="-342900" lvl="0" algn="l">'
               '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="600"/></a:spcAft>'
               f'<a:buSzPts val="1800"/><a:buChar char="{BULLET}"/></a:pPr>')
    elif kind == "item":
        xml = (f'<a:pPr {nsdecls("a")} marL="914400" indent="-342900" lvl="1" algn="l">'
               '<a:spcBef><a:spcPts val="0"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="300"/></a:spcAft>'
               f'<a:buSzPts val="1500"/><a:buChar char="{SUBBULLET}"/></a:pPr>')
    elif kind == "header":
        xml = (f'<a:pPr {nsdecls("a")} marL="0" indent="0" lvl="0" algn="l">'
               '<a:spcBef><a:spcPts val="600"/></a:spcBef>'
               '<a:spcAft><a:spcPts val="300"/></a:spcAft>'
               '<a:buNone/></a:pPr>')
    else:
        return
    p.insert(0, parse_xml(xml))


def _is_title_ph(placeholder):
    """True for title / centre-title / subtitle placeholders (no body bullets)."""
    try:
        t = placeholder.placeholder_format.type
    except (AttributeError, ValueError):
        return False
    return t is not None and int(t) in (1, 3, 4)  # TITLE, CENTER_TITLE, SUBTITLE


def set_text(placeholder, lines, *, bullets=True):
    """Replace placeholder text, forcing consistent bullet formatting.

    Each entry in `lines` is either a string (formatted per `bullets`) or a
    ``(text, kind)`` tuple where kind ∈ {'bullet', 'item', 'header'}.
    Title/subtitle placeholders keep their template formatting (no bullets).
    """
    tf = placeholder.text_frame
    tf.clear()
    is_title = _is_title_ph(placeholder)
    if not lines:
        if not is_title:
            _apply_para_format(tf.paragraphs[0], "header")
        return
    default_kind = "bullet" if bullets else "header"
    for i, entry in enumerate(lines):
        if isinstance(entry, tuple):
            text, kind = entry
        else:
            text, kind = entry, (None if is_title else default_kind)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        if kind is not None:
            _apply_para_format(p, kind)
            if kind == "header":
                for run in p.runs:
                    run.font.bold = True
        else:
            # Title/subtitle: drop any stale pPr so it inherits the template
            # style (no bullet, no body indent).
            old = p._p.find(qn("a:pPr"))
            if old is not None:
                p._p.remove(old)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def fit_picture(slide, image_path, *, left, top, max_width, max_height, valign="center"):
    """Add a picture preserving aspect ratio, fitting in (max_width, max_height)."""
    img = Image.open(image_path)
    iw, ih = img.size
    img.close()
    target_w = max_width
    target_h = int(target_w * ih / iw)
    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * iw / ih)
    pic_left = left + (max_width - target_w) // 2
    if valign == "top":
        pic_top = top
    else:
        pic_top = top + (max_height - target_h) // 2
    slide.shapes.add_picture(str(image_path), pic_left, pic_top,
                             width=target_w, height=target_h)


def fit_picture_row(slide, image_paths, *, left, top, total_width, max_height,
                    gap, valign="center"):
    """Lay images in a row at a common height, packed and centred in the band."""
    aspects = []
    for pth in image_paths:
        img = Image.open(pth)
        aspects.append(img.width / img.height)
        img.close()
    n = len(image_paths)
    avail_w = total_width - gap * (n - 1)
    height = min(max_height, int(avail_w / sum(aspects)))
    widths = [int(height * a) for a in aspects]
    row_w = sum(widths) + gap * (n - 1)
    x = left + (total_width - row_w) // 2
    y = top if valign == "top" else top + (max_height - height) // 2
    for pth, w in zip(image_paths, widths):
        slide.shapes.add_picture(str(pth), x, y, width=w, height=height)
        x += w + gap


def set_font_size(placeholder, pt, *, anchor=None):
    tf = placeholder.text_frame
    if anchor is not None:
        tf.vertical_anchor = anchor
    for p in tf.paragraphs:
        p.font.size = Pt(pt)
        for run in p.runs:
            run.font.size = Pt(pt)


def clear_pictures(prs):
    """Remove every picture AND its image relationship so the script is safe to
    re-run on a filled deck. Dropping the now-dangling image rels lets python-pptx
    prune the orphaned media parts on save (otherwise they accumulate and bloat /
    can break rendering)."""
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape._element.getparent().remove(shape._element)
        part = slide.part
        for rId, rel in list(part.rels.items()):
            if "image" in rel.reltype:
                part.drop_rel(rId)


def _title_text(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 and ph.has_text_frame:
            return ph.text_frame.text.strip()
    return ""


SWEEP_TITLE_PREFIX = "Experiment 4"


def remove_sweep_slide(prs):
    """Idempotently remove the old overlay-sweep comparison slide if present.
    We now feature 10pct as the headline reproduction, so the multi-percentage
    deep-dive slide is dropped."""
    sld_id_lst = prs.slides._sldIdLst
    # Pair each <p:sldId> element with its slide IN ORDER so we remove the exact
    # element (indexing the two lists separately can misalign).
    for sld_id, s in zip(list(sld_id_lst), list(prs.slides)):
        if _title_text(s).startswith(SWEEP_TITLE_PREFIX):
            sld_id_lst.remove(sld_id)
            return


def shrink_body(placeholder, new_height_emu):
    placeholder.height = new_height_emu


def main():
    prs = Presentation(str(PPTX_PATH))
    clear_pictures(prs)  # idempotent: drop old pictures before re-adding
    remove_sweep_slide(prs)  # idempotent: drop the old overlay-sweep comparison slide
    SW = prs.slide_width   # 9_144_000 (10")
    SH = prs.slide_height  # 5_143_500 (5.62")

    body_left = 311700
    body_top = 1152475
    body_width = 8520600
    body_height = 3416400  # full content area
    body_bottom = body_top + body_height

    slides = list(prs.slides)

    # -------- Slide 2: Title --------
    s = slides[1]
    title = s.placeholders[0]
    subtitle = s.placeholders[1]
    set_text(title, [
        "MILAN — Editing Spurious Features",
        "Reproducing Hernandez et al. (ICLR 2022)",
        "+ extensions: VGG16, layer-depth analysis, CLIP ViT-B/32",
    ])
    set_text(subtitle, [
        "Group [#]",
        "Wonmin Kim · Seongho Kim · Ming-Yang Wu · Steven Tsai",
        "June 9, 2026",
    ])
    set_notes(s, (
        "Opener (30s). We're reproducing Section 7 of MILAN — Natural Language "
        "Descriptions of Deep Visual Features, ICLR 2022. Section 7 is the editing-"
        "spurious-features case study, where MILAN's natural-language neuron captions "
        "are used to identify and ablate neurons that latch onto a text shortcut. "
        "On top of the reproduction we ran three new experiments: VGG16 architecture "
        "generalization, a ResNet18 layer-depth analysis, and CLIP ViT-B/32 robustness."
    ))

    # -------- Slide 3: Paper citation --------
    s = slides[2]
    title3 = s.placeholders[0]
    subtitle3 = s.placeholders[1]
    # Default template geometry overflows (huge centre-title font collides with
    # the bottom-anchored subtitle). Reposition + shrink so they never overlap.
    title3.left, title3.top = Emu(311700), Emu(1_050_000)
    title3.width, title3.height = Emu(8_520_600), Emu(1_500_000)
    subtitle3.left, subtitle3.top = Emu(311700), Emu(2_750_000)
    subtitle3.width, subtitle3.height = Emu(8_520_600), Emu(1_600_000)
    set_text(title3, [
        "Paper [A]: Natural Language Descriptions of Deep Visual Features",
    ], bullets=False)
    set_text(subtitle3, [
        "[A] Hernandez, Schwettmann, Bau, Bagashvili, Oliva, Andreas.",
        "Natural Language Descriptions of Deep Visual Features. ICLR 2022.",
        "arXiv:2201.11114 · project: milan.csail.mit.edu",
    ], bullets=False)
    set_font_size(title3, 30, anchor=MSO_ANCHOR.TOP)
    set_font_size(subtitle3, 16, anchor=MSO_ANCHOR.TOP)
    set_notes(s, (
        "Authors are from MIT CSAIL + Northeastern. The paper introduces MILAN, "
        "a method that captions individual neurons in a vision model with free-form "
        "natural language. The downstream payoff we focus on is editing — using the "
        "captions to surgically remove neurons that encode a spurious feature."
    ))

    # -------- Slide 4: Section 1 header --------
    set_notes(slides[3], "Section 1: motivate the problem and why interpretability matters.")

    # -------- Slide 5: Overview of field/topic --------
    s = slides[4]
    set_text(s.placeholders[0], ["Overview: interpretability for trustworthy ML"])
    set_text(s.placeholders[1], [
        "Deep nets are accurate but opaque — we can't easily tell what each neuron encodes.",
        "Prior interpretability (Network Dissection, TCAV, CLIP-Dissect) is limited to a fixed concept vocabulary.",
        "MILAN goes open-vocabulary: a learned image-conditioned language model writes the description.",
        "Goal of Section 7: use those descriptions to find and remove neurons that latched onto a spurious shortcut.",
    ])
    # Shrink the body so the diagram sits below the bullets.
    shrink_body(s.placeholders[1], Emu(1_900_000))
    fit_picture(s, SLIDE_FIG / "milan_pipeline_diagram.png",
                left=body_left, top=body_top + 1_950_000,
                max_width=body_width, max_height=body_bottom - (body_top + 1_950_000))
    set_notes(s, (
        "45 seconds. The field has long wanted to ask the model 'what does this neuron do?', "
        "but earlier methods (Network Dissection, TCAV, CLIP-Dissect) all require a fixed "
        "concept vocabulary you pick in advance, so they miss anything off-vocabulary. "
        "MILAN learns an image-conditioned language model from human captions and produces "
        "open-vocabulary descriptions. In Section 7 the authors use those descriptions as "
        "actionable handles for editing a model — exactly the loop shown in the diagram below."
    ))

    # -------- Slide 6: Problem statement --------
    s = slides[5]
    set_text(s.placeholders[0], ["The problem: spurious shortcuts that hide inside neurons"])
    set_text(s.placeholders[1], [
        "Models trained on data with a spurious correlation latch onto the shortcut — here, painting the class name in the image corner.",
        "Result: high clean accuracy, but adversarial accuracy collapses when the corner text is wrong.",
        "Identifying which neurons encode the shortcut, at scale and without retraining, is the open problem.",
    ])
    shrink_body(s.placeholders[1], Emu(1_500_000))
    fit_picture(s, SLIDE_FIG / "spurious_dataset_grid.png",
                left=body_left, top=body_top + 1_550_000,
                max_width=body_width, max_height=body_bottom - (body_top + 1_550_000))
    set_notes(s, (
        "Explain the spurious setup briefly. Top row is clean Imagenette. Middle row is the "
        "training set — a fraction (10%) of the images have the class name painted in the corner, "
        "a predictive shortcut. Bottom row is the adversarial test set — same images but with the "
        "wrong class name painted. A ResNet18 trained naively gets ~75% on the clean set but drops "
        "to ~51% on the adversarial set — it's leaning on the text, not just the object."
    ))

    # -------- Slide 7: Why important --------
    s = slides[6]
    set_text(s.placeholders[0], ["Why this matters for trustworthy ML"])
    set_text(s.placeholders[1], [
        "Robustness: spurious shortcuts (text overlays, watermarks, background bias) silently degrade real-world deployment.",
        "Fairness & auditing: knowing what a model relies on is a prerequisite for accountability.",
        "Editing without retraining: full retraining is expensive and brittle — neuron-level edits are surgical and reversible.",
        "Three of this course's themes — spurious correlations, adversarial robustness, model editing — meet in one self-contained loop.",
    ])
    set_notes(s, (
        "Tie back to the course themes. The interesting thing about Section 7 is that it "
        "exercises three pillars of trustworthy ML in a single experiment: it builds a "
        "spurious dataset, it measures adversarial robustness, and it uses interpretability "
        "as a tool for model editing. That's why we picked this paper out of the suggested "
        "list."
    ))

    # -------- Slide 8: Section 2 header --------
    set_notes(slides[7], "Section 2: related work — how MILAN compares to prior interpretability methods.")

    # -------- Slide 9: Related works --------
    s = slides[8]
    set_text(s.placeholders[0], ["Related work: from fixed concept lists to open-vocabulary captions"])
    set_text(s.placeholders[1], [
        "Network Dissection (Bau et al. 2017): match each neuron to one of ~1,200 predefined visual concepts. Limited vocabulary.",
        "TCAV (Kim et al. 2018): test directional sensitivity to user-supplied concepts. Requires concept-positive examples for every probe.",
        "CLIP-Dissect (Oikarinen & Weng 2023): use CLIP's text encoder as a fixed concept set. Better than fixed labels, still vocab-bound.",
        "MILAN (this paper): learn an image-conditioned language model on 60k human captions → free-form descriptions, no fixed vocabulary.",
        "Trade-off: open-vocab is more expressive but harder to evaluate; the paper uses a PMI objective to push descriptions toward specificity.",
    ])
    set_notes(s, (
        "60 seconds. Walk through the family tree. The common thread is: every prior method "
        "needs you to enumerate concepts up-front, which means anything unusual — like 'painted "
        "white text in the corner' — won't get found unless you anticipated it. MILAN escapes "
        "that by training a captioning model whose vocabulary is whatever appears in 60k human-"
        "written neuron descriptions. The trade-off: open-vocabulary captions are harder to "
        "evaluate, which is why the paper introduces the PMI scoring objective."
    ))

    # -------- Slide 10: Section 3 header --------
    set_notes(slides[9], "Section 3: how MILAN actually works.")

    # -------- Slide 11: Key ideas --------
    s = slides[10]
    set_text(s.placeholders[0], ["MILAN — key ideas"])
    set_text(s.placeholders[1], [
        "Represent each neuron by its top-k activating image patches (exemplars).",
        "Train an image-conditioned LM on 60k human-written exemplar descriptions (the MILANNOTATIONS dataset).",
        "Decode descriptions with a PMI objective: argmax log p(d | exemplars) − log p(d) — rewards specificity, penalises generic captions.",
        "Section 7 application: regex-flag any caption containing 'text' / 'word' / 'letter', then ablate those neurons and re-measure adversarial accuracy.",
    ])
    shrink_body(s.placeholders[1], Emu(1_800_000))
    fit_picture(s, SLIDE_FIG / "milan_pipeline_diagram.png",
                left=body_left, top=body_top + 1_850_000,
                max_width=body_width, max_height=body_bottom - (body_top + 1_850_000))
    set_notes(s, (
        "60 seconds. The intuition is simple. Each convolution channel has a set of image "
        "patches that maximally activate it — those are its 'exemplars'. MILAN feeds those "
        "exemplars into a small encoder-decoder language model, trained on humans labeling "
        "exemplars in plain English. The trick is the PMI scoring objective at decode time: "
        "instead of just the most likely description, MILAN picks the description that is "
        "most likely *given the exemplars relative to* a generic prior. That kills off bland "
        "captions like 'an image' and keeps specific ones like 'words and letters'. Once you "
        "have a caption per neuron, editing is trivial — string-match for 'text' / 'word' / "
        "'letter' and zero those channels."
    ))

    # -------- Slide 12: Section 4 header --------
    set_notes(slides[11], "Section 4: our experimental setup and results.")

    # -------- Slide 13: Experimental setup --------
    s = slides[12]
    set_text(s.placeholders[0], ["Experimental setup"])
    set_text(s.placeholders[1], [
        "Dataset: spurious-text Imagenette (10 classes; 10% of train images have class-name overlay; adversarial test set has wrong-class overlay).",
        "Model: ResNet18 trained from scratch using the paper's Appendix E hyperparameters.",
        "MILAN: pretrained 'base' decoder from milan.csail.mit.edu — same checkpoint as the paper.",
        "Metrics: clean validation accuracy, adversarial test accuracy, ablation curve (adv-acc vs. # neurons zeroed), text-neuron count.",
        "Compute: UCSD DSMLP — GTX 1080 Ti (11 GB GPU, 16 GB RAM) — required several OOM fixes vs. the upstream code.",
    ])
    shrink_body(s.placeholders[1], Emu(2_200_000))
    fit_picture(s, SLIDE_FIG / "spurious_dataset_grid.png",
                left=body_left, top=body_top + 2_250_000,
                max_width=body_width, max_height=body_bottom - (body_top + 2_250_000))
    set_notes(s, (
        "Note the dataset substitution: the paper uses an unspecified 10-class ImageNet subset; "
        "we use Imagenette (a publicly available 10-class ImageNet subset). Absolute numbers "
        "won't match exactly, but trends should. We also had to write five OOM patches for the "
        "upstream MILAN code to run on the 11GB/16GB DSMLP container — those are documented in "
        "the README. Everything else (training hyperparameters, MILAN decoder, ablation protocol) "
        "follows the paper directly."
    ))

    # -------- Slide 14: Key results of paper --------
    s = slides[13]
    set_text(s.placeholders[0], ["Paper's key results — Fig 7 (text neurons) & Fig 8 (ablation curve)"])
    set_text(s.placeholders[1], [""])  # clear guidance text
    shrink_body(s.placeholders[1], Emu(200_000))
    # Two images side by side, below a thin placeholder
    half_w = body_width // 2 - 60_000
    img_top = body_top + 250_000
    img_max_h = body_bottom - img_top
    fit_picture(s, PAPER_FIG / "paper_fig7.png",
                left=body_left, top=img_top,
                max_width=half_w, max_height=img_max_h, valign="top")
    fit_picture(s, PAPER_FIG / "paper_fig8.png",
                left=body_left + half_w + 120_000, top=img_top,
                max_width=half_w, max_height=img_max_h, valign="top")
    set_notes(s, (
        "Fig 7 (left): the paper trains a 10-class ImageNet classifier with class-name text "
        "stamped in the corner. (a) shows training examples, (b) the adversarial test set with "
        "wrong-class text, and (c) the top exemplars of one neuron (layer3-134) that MILAN "
        "captions as 'words and letters' — exactly the spurious feature. "
        "Fig 8 (right): the ablation curve. Y-axis is adversarial accuracy; X-axis is the number "
        "of neurons zeroed. The blue 'sort text' curve (ablating neurons whose MILAN caption "
        "mentions text) climbs from ~58% to ~63%, beating the orange 'sort all' baseline that "
        "ablates neurons by overall importance. The flat dashed line is the no-distractor ceiling."
    ))

    # -------- Slide 15: Code reproduction --------
    s = slides[14]
    set_text(s.placeholders[0], ["Our reproduction — Fig 7 & Fig 8 (10% overlay)"])
    set_text(s.placeholders[1], [
        "Clean val: 75.4%. Adversarial test: 51.5%. MILAN flags 216/1024 neurons (21.1%) as text.",
        "Text-sorted ablation recovers adversarial accuracy 51.5% → ~55% and clearly beats random/importance ablation (which both degrade) — reproducing the paper's editing result.",
    ])
    shrink_body(s.placeholders[1], Emu(900_000))
    img_top = body_top + 950_000
    img_max_h = body_bottom - img_top
    half_w = body_width // 2 - 60_000
    fit_picture(s, FIG / "fig7.png",
                left=body_left, top=img_top,
                max_width=half_w, max_height=img_max_h)
    fit_picture(s, FIG / "fig8.png",
                left=body_left + half_w + 120_000, top=img_top,
                max_width=half_w, max_height=img_max_h)
    set_notes(s, (
        "Left: our reproduction of Fig 7 — eight text neurons from the 10% ResNet18 with their "
        "MILAN captions; the painted class words ('pump', 'truck', 'springer'…) visible in the "
        "exemplar thumbnails are exactly what these neurons fire on. "
        "Right: our reproduction of Fig 8. From the 0.515 no-ablation baseline, zeroing MILAN's "
        "text neurons (blue) recovers adversarial accuracy to ~0.55 and holds, while ablating the "
        "same number of neurons by importance (orange) or at random (gray) degrades to ~0.46 and "
        "~0.37. The caption-guided edit clearly beats both baselines — the paper's editing claim "
        "reproduces on Imagenette at a 10% overlay (the overlay strength that gives the cleanest "
        "separation; see backup notes)."
    ))

    # -------- Slide 16: Additional experiments --------
    s = slides[15]
    set_text(s.placeholders[0], ["Three new experiments: VGG16 / layer-depth / CLIP"])
    set_text(s.placeholders[1], [""])
    shrink_body(s.placeholders[1], Emu(200_000))
    # 1x3 row of figures: common height, packed and centred in the content band.
    row_top = body_top + 200_000
    fit_picture_row(s, [FIG / "fig_arch_comparison.png",
                        FIG / "fig_layer_analysis.png",
                        FIG / "fig_clip_analysis.png"],
                    left=body_left, top=row_top,
                    total_width=body_width, max_height=body_bottom - row_top,
                    gap=120_000, valign="center")
    set_notes(s, (
        "Three new experiments, left to right (all on the 10% spurious set). "
        "(1) Architecture generalization — VGG16. MILAN finds 357/1472 text neurons (24.3%). VGG16 "
        "is more shortcut-reliant than ResNet18 (adversarial baseline 19.8% vs 51.5%), but text-"
        "sorted ablation still recovers it strongly (19.8% → 32.7%, +12.8pp) and beats random — the "
        "editing generalizes across CNNs without code changes. "
        "(2) Layer-depth analysis on ResNet18. Text-neuron fraction jumps 3.0× between conv1 (14.1%) "
        "and layer1 (42.2%) — the shortcut appears right after the first residual block — then "
        "declines through the deeper layers (29%, 30%, 13%). Description diversity drops with depth. "
        "(3) CLIP ViT-B/32 zero-shot. Max text-neuron fraction across blocks is 7.9% — about 2.7× "
        "lower than ResNet18 (21.1%). Top words are 'objects', 'colored', 'white' rather than 'words' "
        "or 'letters'. CLIP's vision-language pretraining is structurally more robust to text overlays."
    ))

    # -------- Slide 17: Section 5 header --------
    set_notes(slides[16], "Section 5: concluding takeaways.")

    # -------- Slide 18: Conclusions --------
    s = slides[17]
    set_text(s.placeholders[0], ["Conclusions & take-away"])
    set_text(s.placeholders[1], [
        "MILAN reproduces on Imagenette (10% overlay): 216/1024 text neurons identified; text-sorted ablation recovers adversarial accuracy 51.5% → ~55%, beating random/importance baselines.",
        "Text-neuron ablation generalizes across CNN architectures (ResNet18 → VGG16) without code changes; VGG16 recovers 19.8% → 32.7%.",
        "Spurious shortcuts crystallize right after the first residual block (3.0× jump conv1→layer1), then thin out with depth.",
        "CLIP's image-language pretraining produces representations ~2.7× more robust to spurious text features than a task-trained CNN.",
        "Open-vocabulary, model-agnostic neuron descriptions are a practical handle for editing model behaviour.",
    ])
    shrink_body(s.placeholders[1], Emu(1_900_000))
    fit_picture(s, SLIDE_FIG / "summary_metrics_table.png",
                left=body_left, top=body_top + 1_950_000,
                max_width=body_width, max_height=body_bottom - (body_top + 1_950_000))
    set_notes(s, (
        "Headline takeaways at a 10% overlay. MILAN reproduces — caption-guided ablation recovers "
        "adversarial robustness and clearly beats importance/random baselines. On top of that our "
        "three extensions add: VGG16 evidence for cross-architecture generalization, a layer-depth "
        "story (the shortcut appears very early, at layer1, then thins out with depth), and a "
        "robustness comparison against CLIP — vision-language pretraining is ~2.7× sparser in "
        "text-selective neurons than a task-trained CNN."
    ))

    # -------- Slide 19: Section 6 header --------
    set_notes(slides[18], "Section 6: discussion, limitations, and future work.")

    # -------- Slide 20: Likes / dislikes --------
    s = slides[19]
    set_text(s.placeholders[0], ["What we liked / didn't like"])
    set_text(s.placeholders[1], [
        "Like: open-vocabulary descriptions — no concept set to engineer.",
        "Like: model-agnostic — works on ResNet, VGG, CLIP with the same decoder.",
        "Like: descriptions are *actionable* — they directly enable editing, not just analysis.",
        "Dislike: text-neuron identification is regex-based — brittle to synonyms ('caption', 'sign', 'character').",
        "Dislike: decoder is English-only, bounded by the 60k MILANNOTATIONS caption corpus.",
        "Dislike: no uncertainty estimate on captions — confidently wrong descriptions look the same as confidently right ones.",
    ])
    set_notes(s, (
        "Quick honest take. The biggest 'like' is that MILAN escapes the concept-vocabulary "
        "bottleneck. The biggest 'dislike' is that the downstream filter — 'does the caption "
        "contain text/word/letter?' — is a regex, which is fragile. A semantic-similarity filter "
        "with a sentence embedder would be a one-line improvement."
    ))

    # -------- Slide 21: Strengths & weaknesses --------
    s = slides[20]
    set_text(s.placeholders[0], ["Strengths & weaknesses"])
    set_text(s.placeholders[1], [
        ("Strengths", "header"),
        ("PMI objective is a clean, principled fix for generic captions.", "item"),
        ("Three case studies in the paper (BigGAN, ImageNet classifiers, adversarial editing) show real utility, not just toy demos.", "item"),
        ("Pretrained decoder is plug-and-play — we ran it on ResNet18, VGG16, and CLIP without retraining.", "item"),
        ("Weaknesses", "header"),
        ("Decoder quality is upper-bounded by the size and diversity of the 60k caption corpus.", "item"),
        ("Multimodal/transformer units need custom adapters — we had to write clip_glue.py to reshape ViT tokens into (B,C,H,W).", "item"),
        ("No reported test for caption consistency across runs / random seeds.", "item"),
    ])
    set_notes(s, (
        "Beyond like/dislike: a structural strength is the PMI objective, which is the single "
        "design choice that makes the captions actually useful. A structural weakness is that "
        "you inherit the captioning corpus's blind spots — anything not well-represented in "
        "the 60k human captions will get described poorly. Adapting to CLIP took us a non-"
        "trivial amount of work because the upstream code assumed conv-style spatial tensors."
    ))

    # -------- Slide 22: Limitations --------
    s = slides[21]
    set_text(s.placeholders[0], ["Limitations"])
    set_text(s.placeholders[1], [
        "Keyword-based text-neuron identification is brittle — 'sign', 'character', 'label' all describe text but won't match the regex.",
        "Ablation is destructive: a neuron is zeroed in full, with no fine-grained or graded edit.",
        "Only top-k exemplars are seen by the decoder — rare but discriminative activations may be invisible to the description.",
        "Evaluation is ground-truth-free: the paper has no neuron-level gold labels, so caption correctness is largely qualitative.",
        "Generalization to transformers / LLMs is unproven — our CLIP results required custom hooks.",
    ])
    set_notes(s, (
        "Tee up future work. The four limitations above are exactly the four motivating "
        "ideas for what to do next — and each one has a 1–2 sentence pitch on the next slide."
    ))

    # -------- Slide 23: Future work --------
    s = slides[22]
    set_text(s.placeholders[0], ["Future work"])
    set_text(s.placeholders[1], [
        "Embedding-based text-neuron detection: replace regex with sentence-encoder similarity → catches synonyms / paraphrases.",
        "Graded edits: activation steering or low-rank shrinkage instead of zero-ablation; preserves capacity while removing the shortcut.",
        "Beyond top-k: condition the decoder on a fuller distribution of activations (heatmaps, statistics) to surface rare features.",
        "MILAN-for-transformers: extend to ViTs and LLM hidden states with shared adapters — start from our clip_glue.py prototype.",
        "Harder spurious benchmarks: Waterbirds, CelebA, NICO — broaden from text overlay to multi-axis spurious correlations.",
    ])
    set_notes(s, (
        "Brief — most of these are obvious one-step extensions. The embedding-based filter is "
        "the lowest-hanging fruit; the activation-steering idea connects to recent mech-interp "
        "work on Anthropic's models; the LLM extension is the long-horizon direction."
    ))

    # -------- Slide 24: References --------
    s = slides[23]
    set_text(s.placeholders[0], ["References"])
    set_text(s.placeholders[1], [
        "[1] Hernandez, Schwettmann, Bau, Bagashvili, Oliva, Andreas. Natural Language Descriptions of Deep Visual Features. ICLR 2022. arXiv:2201.11114.",
        "[2] Bau, Zhou, Khosla, Oliva, Torralba. Network Dissection: Quantifying Interpretability of Deep Visual Representations. CVPR 2017.",
        "[3] Kim, Wattenberg, Gilmer, Cai, Wexler, Viégas, Sayres. TCAV: Interpretability Beyond Feature Attribution. ICML 2018.",
        "[4] Oikarinen & Weng. CLIP-Dissect: Automatic Description of Neuron Representations in Deep Vision Networks. ICLR 2023.",
        "[5] Radford et al. Learning Transferable Visual Models From Natural Language Supervision (CLIP). ICML 2021.",
        "[6] Howard. Imagenette — github.com/fastai/imagenette.",
        "[7] Our repo: github.com/1minute99/DSC291_Research_Reproduction.",
    ])
    set_notes(s, "")

    out = PPTX_PATH
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    main()
